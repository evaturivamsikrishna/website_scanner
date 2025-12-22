from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()

    # Slide 1: Title Slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Automated Link Checking vs. Manual Verification"
    subtitle.text = "Efficiency, Accuracy, and Scalability for Kwalee.com"

    # Slide 2: The Manual Process
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "The Manual Challenge"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Checking links manually is a bottleneck:"
    
    p = tf.add_paragraph()
    p.text = "• Time-consuming: Checking 1000+ links takes ~16 hours."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Error-prone: Human fatigue leads to missed broken links."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Inconsistent: Difficult to maintain a twice-daily schedule."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• No Visibility: No centralized data or historical trends."
    p.level = 1

    # Slide 3: The Automated Solution
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "The Automated Solution"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Our new system transforms the workflow:"
    
    p = tf.add_paragraph()
    p.text = "• Fast: Async checking handles thousands of links in minutes."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Scheduled: Runs automatically 2x daily via GitHub Actions."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Comprehensive: Deep English crawl + global locale mapping."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Centralized: Real-time dark-themed dashboard with charts."
    p.level = 1

    # Slide 4: Time Comparison
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Time Comparison: Manual vs. Automation"
    
    # Add a table
    rows, cols = 3, 3
    left, top, width, height = Inches(1), Inches(2), Inches(8), Inches(2)
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Set column widths
    table.columns[0].width = Inches(3)
    table.columns[1].width = Inches(2.5)
    table.columns[2].width = Inches(2.5)
    
    # Write headers
    table.cell(0, 0).text = "Metric"
    table.cell(0, 1).text = "Manual Check"
    table.cell(0, 2).text = "Automation"
    
    # Row 1
    table.cell(1, 0).text = "Time (1000 Links)"
    table.cell(1, 1).text = "~16.6 Hours"
    table.cell(1, 2).text = "~5 Minutes"
    
    # Row 2
    table.cell(2, 0).text = "Frequency"
    table.cell(2, 1).text = "Occasional"
    table.cell(2, 2).text = "Twice Daily"

    # Slide 5: Key Benefits
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Key Benefits & ROI"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    p = tf.add_paragraph()
    p.text = "• 99.5% Reduction in verification time."
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• Instant Alerts: Email notifications when links break."
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• SEO Protection: Prevent 'Page Not Found' errors for users."
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• Resource Optimization: Team focuses on fixing, not finding."
    p.level = 0

    # Slide 6: Conclusion
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Conclusion"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Automation provides a scalable, reliable, and efficient way to maintain Kwalee's web presence."

    # Save the presentation
    prs.save('Link_Checker_Comparison.pptx')
    print("✅ Presentation created: Link_Checker_Comparison.pptx")

if __name__ == "__main__":
    create_presentation()
